#!/usr/bin/env python3
"""Generates docs/user-guide.pptx for the Morph Ops Control Room."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- design tokens (match the HTML guide) ----
PRIMARY      = RGBColor(0xFF, 0x91, 0x4D)
PRIMARY_DARK = RGBColor(0xE5, 0x82, 0x2D)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1A, 0x1A, 0x2E)
MUTED        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG     = RGBColor(0xFF, 0xF0, 0xE6)
BORDER       = RGBColor(0xE5, 0xE7, 0xEB)
ADMIN        = RGBColor(0xF5, 0x9E, 0x0B)
FACILITATOR  = RGBColor(0x7C, 0x3A, 0xED)
CM           = RGBColor(0x16, 0xA3, 0x4A)
ROW_ALT      = RGBColor(0xFA, 0xF6, 0xF1)
GREEN_TICK   = RGBColor(0x16, 0xA3, 0x4A)
GREY_DASH    = RGBColor(0xB0, 0xB0, 0xC0)

SW, SH = Inches(13.333), Inches(7.5)
FONT = "Poppins"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _box(slide, l, t, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)  # rectangle
    shp.fill.solid() if fill else shp.fill.background()
    if fill:
        shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set(r, size, color, bold)
    return tb


def slide_cover(prs):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=PRIMARY)
    card_w, card_h = Inches(9), Inches(3.6)
    cl = (SW - card_w) / 2
    ct = (SH - card_h) / 2
    _box(s, cl, ct, card_w, card_h, fill=WHITE)
    _text(s, cl, ct + Inches(0.6), card_w, Inches(0.5), "MORPH BY TNX", 16, PRIMARY_DARK, True, PP_ALIGN.CENTER)
    _text(s, cl, ct + Inches(1.15), card_w, Inches(1.0), "Ops Control Room", 40, DARK, True, PP_ALIGN.CENTER)
    _text(s, cl, ct + Inches(2.15), card_w, Inches(0.5), "User Guide", 20, MUTED, False, PP_ALIGN.CENTER)
    _text(s, cl, ct + Inches(2.75), card_w, Inches(0.5), "Version 1.0  ·  May 2026", 13, MUTED, False, PP_ALIGN.CENTER)
    return s


def slide_section(prs, num, title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _box(s, 0, 0, Inches(0.35), SH, fill=PRIMARY)
    _text(s, Inches(7.5), Inches(0.2), Inches(5.5), Inches(4), str(num), 200, LIGHT_BG, True, PP_ALIGN.RIGHT)
    _text(s, Inches(1.0), Inches(2.7), Inches(10), Inches(1.2), title, 40, DARK, True)
    if subtitle:
        _text(s, Inches(1.05), Inches(3.8), Inches(10), Inches(0.8), subtitle, 18, MUTED)
    return s


def slide_content(prs, title, bullets, note=""):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _box(s, 0, 0, SW, Inches(0.18), fill=PRIMARY)        # top bar
    _box(s, 0, 0, Inches(0.18), SH, fill=PRIMARY)        # left strip
    _box(s, Inches(0.18), Inches(0.18), SW - Inches(0.18), Inches(1.1), fill=LIGHT_BG)  # header band
    _text(s, Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.8), title, 26, DARK, True, anchor=MSO_ANCHOR.MIDDLE)

    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if b.startswith("## "):
            r = p.add_run(); r.text = b[3:]
            _set(r, 18, PRIMARY_DARK, True)
            p.space_before = Pt(10)
        elif b.startswith("-- "):
            p.level = 1
            r = p.add_run(); r.text = "·  " + b[3:]
            _set(r, 15, MUTED)
        else:
            r = p.add_run(); r.text = "•  " + b
            _set(r, 17, DARK)
            p.space_after = Pt(6)
    if note:
        nb = _box(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.85), fill=LIGHT_BG)
        ntf = nb.text_frame; ntf.word_wrap = True
        ntf.margin_left = Inches(0.2); ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = ntf.paragraphs[0].add_run(); r.text = "Tip:  " + note
        _set(r, 13, PRIMARY_DARK, True)
    return s


def slide_roles(prs, roles):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _box(s, 0, 0, SW, Inches(0.18), fill=PRIMARY)
    _text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8), "Roles & Access", 30, DARK, True)
    card_w = Inches(3.85); gap = Inches(0.3)
    total = card_w * 3 + gap * 2
    start = (SW - total) / 2
    top = Inches(1.9); card_h = Inches(4.6)
    for i, (name, color, can, cannot) in enumerate(roles):
        l = start + i * (card_w + gap)
        _box(s, l, top, card_w, card_h, fill=WHITE, line=BORDER)
        _box(s, l, top, card_w, Inches(0.85), fill=color)
        _text(s, l, top, card_w, Inches(0.85), name, 18, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        tb = s.shapes.add_textbox(l + Inches(0.25), top + Inches(1.0), card_w - Inches(0.5), card_h - Inches(1.2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; r = p.add_run(); r.text = "CAN"; _set(r, 12, color, True)
        for c in can:
            p = tf.add_paragraph(); r = p.add_run(); r.text = "• " + c; _set(r, 12.5, DARK)
        p = tf.add_paragraph(); p.space_before = Pt(8); r = p.add_run(); r.text = "CANNOT"; _set(r, 12, MUTED, True)
        for c in cannot:
            p = tf.add_paragraph(); r = p.add_run(); r.text = "• " + c; _set(r, 12.5, MUTED)
    return s


def slide_table(prs, title, headers, rows):
    s = prs.slides.add_slide(BLANK)
    _box(s, 0, 0, SW, SH, fill=WHITE)
    _box(s, 0, 0, SW, Inches(0.18), fill=PRIMARY)
    _text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8), title, 30, DARK, True)
    rows_n = len(rows) + 1
    tbl_w = Inches(11.5); tbl_h = Inches(0.6) * rows_n
    gx = s.shapes.add_table(rows_n, len(headers), Inches(0.9), Inches(1.8), tbl_w, tbl_h).table
    for c, h in enumerate(headers):
        cell = gx.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = h; _set(r, 13, WHITE, True)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = gx.cell(ri, ci); cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if ri % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            if val == "Y":
                r = p.add_run(); r.text = "✓"; _set(r, 14, GREEN_TICK, True); p.alignment = PP_ALIGN.CENTER
            elif val == "-":
                r = p.add_run(); r.text = "—"; _set(r, 14, GREY_DASH); p.alignment = PP_ALIGN.CENTER
            else:
                r = p.add_run(); r.text = val; _set(r, 12.5, DARK, ci == 0)
    return s


# ============ BUILD DECK ============
slide_cover(prs)

slide_section(prs, 1, "Getting Started", "Install the app and sign in")
slide_content(prs, "Install & First Sign-in", [
    "## Install (PWA)",
    "iPhone (Safari): Share → Add to Home Screen",
    "Android (Chrome): menu → Add to Home screen / Install app",
    "Allow notifications when prompted — that's how reminders reach you",
    "## First sign-in",
    "Open your login link, sign in with the temporary password",
    "Set a new password once (complete-setup)",
    "You land on the Dashboard; your menu shows only what your role allows",
], note="The app is the source of truth — if it isn't in the app, it didn't happen.")

slide_section(prs, 2, "Roles & Access", "Three roles decide what you see and edit")
slide_roles(prs, [
    ("Admin", ADMIN, ["Everything", "Settings & users", "Cohorts, reviews, reminders", "Import & export"], ["—"]),
    ("Facilitator", FACILITATOR, ["View all data", "Action reviews", "Action sessions", "Dashboards & alerts"], ["Manage users", "Change settings", "Import / export"]),
    ("Community Mgr", CM, ["See every page", "Edit Participants", "Edit Community", "Work tasks, flag risk", "Raise escalations"], ["Edit other modules", "Manage users/settings"]),
])
slide_table(prs, "Permissions at a glance", ["Capability", "Admin", "Facilitator", "CM"], [
    ["View all pages", "Y", "Y", "Y"],
    ["Edit Participants", "Y", "-", "Y"],
    ["Edit Community", "Y", "-", "Y"],
    ["Action Reviews", "Y", "Y", "-"],
    ["Settings & Users", "Y", "-", "-"],
    ["Import / Export", "Y", "-", "-"],
])

slide_section(prs, 3, "Participants & Risk", "The roster — and how to flag risk early")
slide_content(prs, "Working the Participants page", [
    "Open Participants and click a student's row",
    "Update attendance & submission the SAME day it happens",
    "Add a Next action so the next person knows the plan",
    "## Flagging risk",
    "Set Risk = Red (or Amber) early, not at week's end",
], note="Risk = Red auto-creates an outreach task on the owning CM — flagging is how you trigger help.")

slide_section(prs, 4, "Reviews, Sessions & Community", "The weekly operating rhythm")
slide_content(prs, "Assignment Reviews", [
    "Submissions are grouped per week in the Reviews queue",
    "Open a review card — submission details show inline",
    "Log feedback against the correct week",
    "-- Toggle a week open/closed as needed",
    "-- Admins can copy a submission link from Reviews settings",
])
slide_content(prs, "Sessions & Attendance", [
    "Each session has an editable readiness checklist",
    "Share the public attendance link (/attendance/…)",
    "Participants sign in & out with no account",
    "Attendance flows back into the session record",
])
slide_content(prs, "Community (CM Tracker)", [
    "File your weekly row at end of week (Friday)",
    "## Capture",
    "Prompts posted, attendance/submissions updated",
    "Silent count, stuck count, escalations raised",
    "Energy level, key concerns, next actions",
    "Set Status = Done and Weekly report sent = yes",
])

slide_section(prs, 5, "Tasks & Escalations", "Accountability and safeguarding")
slide_content(prs, "Your Tasks queue", [
    "Tasks are owned by role labels (e.g. 'CM Owner')",
    "Change status inline; completed tasks hide by default",
    "Comment what's blocked — @mentions render as chips",
    "Recurring duties arrive automatically and nudge you",
], note="Don't let an overdue task sit silently — comment the blocker.")
slide_content(prs, "Safeguarding Escalations", [
    "Anything you can't resolve goes through escalations — not DMs",
    "Bump 'Escalations raised' on your weekly Community row",
    "Write the detail in 'Key concerns'",
    "It surfaces on the Dashboard for the core team",
])

slide_section(prs, 6, "Notifications & Admin", "Reminders, import and backups")
slide_content(prs, "Notifications & Reminders", [
    "Allow web push when installing the PWA",
    "Daily reminder at 08:00 for your role's due/overdue tasks",
    "The Alerts page collects in-app notifications",
])
slide_content(prs, "Admin: Import & Export", [
    "## Import (one-time)",
    "Go to /admin/import, type IMPORT_MORPH_OPS",
    "Upload the workbook, map columns — records land in the right modules",
    "## Backups",
    "Export JSON from /admin/export regularly",
], note="First Supabase user signs in as admin when none exists (claim_first_admin()).")

slide_section(prs, 7, "FAQ", "Common questions")
slide_content(prs, "Frequently Asked Questions", [
    "## Don't see a page?",
    "-- Pages are role-gated; ask an admin if you need access",
    "## Flagged Red but no task?",
    "-- It auto-creates one on the owning CM — check their Tasks",
    "## Submission in the wrong week?",
    "-- It lands in the bucket of the submit link used",
    "## No reminders?",
    "-- Install the PWA and allow notifications",
    "## Do participants log in?",
    "-- No — public submit & attendance links only",
])

prs.save("user-guide.pptx")
print("Saved user-guide.pptx")
